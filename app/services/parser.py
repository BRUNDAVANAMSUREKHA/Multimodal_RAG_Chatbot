import os
import unicodedata


def _clean_text(text: str) -> str:
    if not text:
        return text
    text = unicodedata.normalize("NFC", text)
    text = text.encode("utf-8", errors="replace").decode("utf-8", errors="replace")
    return text


def parse_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt":
        text = _parse_txt(file_path)
    elif ext == ".pdf":
        text = _parse_pdf(file_path)
    elif ext == ".docx":
        text = _parse_docx(file_path)
    elif ext == ".csv":
        text = _parse_csv(file_path)
    elif ext == ".pptx":
        text = _parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    return _clean_text(text)


def _parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _parse_pdf(file_path):
    try:
        from pypdf import PdfReader
        reader = PdfReader(file_path)
        pages = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append(f"[Page {i + 1}]\n{text.strip()}")
        return "\n\n".join(pages)
    except Exception as e:
        raise RuntimeError(f"PDF parsing failed: {str(e)}")


def _parse_docx(file_path):
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def _parse_csv(file_path):
    import pandas as pd
    df = pd.read_csv(file_path, encoding="utf-8", encoding_errors="replace")
    return df.to_string(index=False)


def _parse_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"[Slide {i + 1}]\n" + "\n".join(texts))
    return "\n\n".join(slides)